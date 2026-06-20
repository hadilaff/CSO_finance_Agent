"""RAG: parse uploaded docs (PDF/DOCX/PPTX/TXT), chunk, embed with Gemini,
store and search in a persistent Chroma collection."""
from __future__ import annotations

import io
import re
import time
import logging
from pathlib import Path

import chromadb
from chromadb.config import Settings
from chromadb.api.types import Documents, EmbeddingFunction, Embeddings

from config import CHROMA_DIR, EMBED_MODEL, get_embed_client

logger = logging.getLogger(__name__)

# Embedding models to try in order when one is unavailable
EMBED_MODEL_FALLBACKS = [
    EMBED_MODEL,                    # primary: models/gemini-embedding-001
    "models/gemini-embedding-2",    # newer model, fallback
]

COLLECTION_NAME = "institutional_knowledge"


# ---------- Document parsing ----------

def _clean_cell(c) -> str:
    """Normalize a pdfplumber table cell to a single-line string."""
    if c is None:
        return ""
    return re.sub(r"\s+", " ", str(c)).strip()


# Strategies tried in order. First strategy that finds at least one valid
# table wins. "text" strategy is what catches borderless tables (KPI grids
# with no drawn lines) — default "lines" strategy misses them entirely.
_TABLE_STRATEGIES: list[dict | None] = [
    None,  # pdfplumber default (line-based)
    {"vertical_strategy": "text", "horizontal_strategy": "text",
     "intersection_y_tolerance": 10, "intersection_x_tolerance": 10},
    {"vertical_strategy": "text", "horizontal_strategy": "lines"},
]

# Section headings look like "4. Key Performance Indicators" or "1. Initiative…"
_HEADING_RE = re.compile(r"^\s*\d+\.\s+[A-Z]")


def _find_tables(page):
    """Try several strategies; return list of (table_obj, rows) pairs."""
    for settings in _TABLE_STRATEGIES:
        try:
            tables = page.find_tables(table_settings=settings) if settings else page.find_tables()
        except Exception as e:
            logger.debug("find_tables failed with %s: %s", settings, e)
            continue
        good: list[tuple] = []
        for t in tables:
            try:
                rows = t.extract()
            except Exception:
                continue
            # Require header + at least one data row, and 2+ columns
            if rows and len(rows) >= 2 and any(len([c for c in r if c]) >= 2 for r in rows):
                good.append((t, rows))
        if good:
            return good
    return []


def _find_table_title(page, table_bbox) -> str:
    """Return the section heading immediately above a table (closest line within 80pt)."""
    x0, top, x1, bottom = table_bbox
    try:
        words = page.extract_words()
    except Exception:
        return ""
    above = [w for w in words if w["bottom"] <= top and w["bottom"] >= top - 80]
    if not above:
        return ""
    # Cluster words into lines by rounded y-coordinate
    lines: dict[int, list] = {}
    for w in above:
        key = int(round(w["top"]))
        lines.setdefault(key, []).append(w)
    # Iterate from closest line to table upward; first heading-like line wins
    for key in sorted(lines.keys(), reverse=True):
        line_words = sorted(lines[key], key=lambda w: w["x0"])
        text = " ".join(w["text"] for w in line_words).strip()
        if _HEADING_RE.match(text):
            return text
    # No numbered heading found; fall back to the closest non-empty line if short
    for key in sorted(lines.keys(), reverse=True):
        line_words = sorted(lines[key], key=lambda w: w["x0"])
        text = " ".join(w["text"] for w in line_words).strip()
        if 5 <= len(text) <= 80:
            return text
    return ""


def _table_to_blocks(title: str, rows: list[list[str]]) -> str:
    """Render a table as: markdown grid + per-row sentence + per-cell sentence.

    The per-cell sentences are what make queries like "digital asset licenses
    Q2 2026" hit. Each one is a self-contained natural-language fact:
        "In '4. Key Performance Indicators', Digital Asset Licenses Q2 2026: 28."
    The embedding for the question matches that sentence almost directly.
    """
    rows = [[_clean_cell(c) for c in r] for r in rows if r and any(_clean_cell(c) for c in r)]
    if len(rows) < 2:
        return ""
    headers, *body = rows
    n = len(headers)
    headers = [h if h else f"col{i + 1}" for i, h in enumerate(headers)]
    body = [(r + [""] * n)[:n] for r in body]

    title_prefix = f"In '{title}', " if title else ""
    title_label = f" — {title}" if title else ""

    md_lines = [f"[Table{title_label}]",
                "| " + " | ".join(headers) + " |",
                "| " + " | ".join(["---"] * n) + " |"]
    md_lines += ["| " + " | ".join(r) + " |" for r in body]

    sentences: list[str] = []
    for r in body:
        label = r[0] or "row"
        cells = [(h, v) for h, v in zip(headers[1:], r[1:]) if v]
        if not cells:
            continue
        # Whole-row sentence
        row_summary = ", ".join(f"{h} = {v}" for h, v in cells)
        sentences.append(f"{title_prefix}{label}: {row_summary}.")
        # Per-cell sentence: highest-signal unit for retrieval
        for h, v in cells:
            sentences.append(f"{title_prefix}{label} {h}: {v}.")

    block = "\n".join(md_lines)
    if sentences:
        # Use blank line so the chunker prefers to keep these together,
        # but will fall back to single-newline boundary if the block is large.
        block += "\n\n[Table rows in natural language]\n" + "\n".join(sentences)
    return block


def _parse_pdf(data: bytes) -> list[tuple[int, str]]:
    """Return list of (page_number, text) tuples — one per page.

    Per page: extract free text, detect tables (multiple strategies for
    borderless tables), find the section heading above each table, then emit
    markdown grid + natural-language sentences per row/cell. Falls back to
    pypdf if pdfplumber is unavailable (tables will be flattened).
    """
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed; falling back to pypdf (tables will be flattened).")
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return [
            (i, page.extract_text() or "")
            for i, page in enumerate(reader.pages, start=1)
            if (page.extract_text() or "").strip()
        ]

    pages: list[tuple[int, str]] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            table_blocks: list[str] = []
            try:
                for tbl, rows in _find_tables(page):
                    title = _find_table_title(page, tbl.bbox)
                    block = _table_to_blocks(title, rows)
                    if block:
                        table_blocks.append(block)
            except Exception as e:
                logger.warning("Table extraction failed on page %d: %s", i, e)
            parts = [text] if text.strip() else []
            parts.extend(table_blocks)
            combined = "\n\n".join(parts).strip()
            if combined:
                pages.append((i, combined))
            logger.info("Page %d: text=%d chars, %d table(s) extracted",
                        i, len(text), len(table_blocks))
    return pages


def _parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # also pull table text
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n".join(paragraphs)


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    pres = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(pres.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(parts)


def _parse_text(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("latin-1", errors="ignore")


def parse_file(filename: str, data: bytes) -> str:
    """Parse a file and return plain text (page markers embedded for PDFs)."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        pages = _parse_pdf(data)
        return "\n\n".join(f"[Page {p}]\n{t}" for p, t in pages)
    parsers = {
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".txt": _parse_text,
        ".md": _parse_text,
    }
    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return parser(data)


def parse_file_with_pages(filename: str, data: bytes) -> list[tuple[int, str]]:
    """Parse a file and return (page_number, text) pairs.

    For non-PDF formats, the whole document is treated as page 1.
    """
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(data)
    text = parse_file(filename, data)
    return [(1, text)] if text.strip() else []


# ---------- Chunking ----------

def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks at paragraph/sentence boundaries."""
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    chunks: list[str] = []
    i = 0
    n = len(text)
    while i < n:
        end = min(i + chunk_size, n)
        if end < n:
            # prefer to break at a natural boundary in the second half of the window
            for sep in ("\n\n", "\n", ". ", " "):
                idx = text.rfind(sep, i + chunk_size // 2, end)
                if idx > i:
                    end = idx + len(sep)
                    break
        chunk = text[i:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        i = max(end - overlap, i + 1)
    return chunks


# ---------- Gemini embedding function for Chroma ----------

def _embed_with_retry(
    texts: list[str],
    max_retries: int = 5,
    base_delay: float = 2.0,
) -> Embeddings:
    """Embed texts using Gemini with exponential backoff and model fallback.

    Retries on 503 (UNAVAILABLE) and 429 (RESOURCE_EXHAUSTED). After
    ``max_retries`` failures on the primary model, tries fallback models once
    each before raising.
    """
    last_exc: Exception | None = None

    for model in dict.fromkeys(EMBED_MODEL_FALLBACKS):  # deduplicated, order kept
        delay = base_delay
        for attempt in range(1, max_retries + 1):
            try:
                client = get_embed_client()
                result = client.models.embed_content(
                    model=model,
                    contents=texts,
                )
                return [list(e.values) for e in result.embeddings]
            except Exception as e:
                err_str = str(e)
                is_retryable = any(
                    code in err_str for code in ("503", "429", "UNAVAILABLE", "RESOURCE_EXHAUSTED")
                )
                if not is_retryable:
                    raise  # non-transient error — don't retry
                last_exc = e
                wait = delay * (2 ** (attempt - 1))
                logger.warning(
                    "Embedding attempt %d/%d failed for model %s (%s). "
                    "Retrying in %.1fs…",
                    attempt, max_retries, model, err_str[:120], wait,
                )
                time.sleep(wait)

        logger.warning("All %d retries exhausted for model %s. Trying next fallback…", max_retries, model)

    raise RuntimeError(
        f"Gemini embedding failed after retrying all models. Last error: {last_exc}"
    )


class GeminiEmbeddingFunction(EmbeddingFunction):
    """Calls Gemini's embed_content for both indexing and querying,
    with automatic retry + model fallback on transient errors."""

    def __init__(self, model: str = EMBED_MODEL):
        self.model = model

    def __call__(self, input: Documents) -> Embeddings:
        return _embed_with_retry(list(input))


# ---------- Chroma collection ----------

_client: chromadb.PersistentClient | None = None
_collection = None


def _get_collection():
    global _client, _collection
    if _collection is None:
        _client = chromadb.PersistentClient(
            path=str(CHROMA_DIR),
            settings=Settings(anonymized_telemetry=False),
        )
        _collection = _client.get_or_create_collection(
            name=COLLECTION_NAME,
            embedding_function=GeminiEmbeddingFunction(),
            metadata={"hnsw:space": "cosine"},
        )
    return _collection


_TABLE_ROWS_RE = re.compile(
    r"\[Table rows in natural language\]\n(.+?)(?=\n\n|\Z)",
    re.S,
)


def _extract_table_sentences(page_text: str) -> list[str]:
    """Pull each '[Table rows in natural language]' line out as a standalone sentence.

    Indexing each sentence atomically means a query like 'digital asset licenses
    Q2 2026' lands directly on the chunk that contains '... Q2 2026: 28.' rather
    than competing with the surrounding prose chunk for the top-k window.
    """
    sentences: list[str] = []
    for m in _TABLE_ROWS_RE.finditer(page_text):
        for line in m.group(1).split("\n"):
            line = line.strip()
            if line:
                sentences.append(line)
    return sentences


def index_file(filename: str, data: bytes) -> int:
    """Parse, chunk, embed, and store one file. Returns the number of chunks indexed.

    Each page is split into two kinds of chunks:
      - kind='text'      — windowed chunks of the full page text (prose + tables in context)
      - kind='table_row' — one chunk per '[Table rows in natural language]' sentence,
                           so cell-level queries can rank on exact keyword overlap.
    """
    pages = parse_file_with_pages(filename, data)
    if not pages:
        return 0

    all_chunks: list[str] = []
    all_metadatas: list[dict] = []
    chunk_idx = 0
    for page_num, page_text in pages:
        for chunk in chunk_text(page_text):
            all_chunks.append(chunk)
            all_metadatas.append({
                "source": filename,
                "chunk": chunk_idx,
                "page": page_num,
                "kind": "text",
            })
            chunk_idx += 1
        for sentence in _extract_table_sentences(page_text):
            all_chunks.append(sentence)
            all_metadatas.append({
                "source": filename,
                "chunk": chunk_idx,
                "page": page_num,
                "kind": "table_row",
            })
            chunk_idx += 1

    if not all_chunks:
        return 0

    coll = _get_collection()
    # Remove any prior chunks for this filename (re-upload safe).
    try:
        coll.delete(where={"source": filename})
    except Exception:
        pass
    ids = [f"{filename}::{i}" for i in range(len(all_chunks))]
    coll.add(documents=all_chunks, ids=ids, metadatas=all_metadatas)
    return len(all_chunks)


def search(query: str, k: int = 5) -> list[dict]:
    coll = _get_collection()
    n = coll.count()
    if n == 0:
        return []
    result = coll.query(query_texts=[query], n_results=min(k, n))
    hits: list[dict] = []
    for doc, meta, dist in zip(
        result["documents"][0],
        result["metadatas"][0],
        result["distances"][0],
    ):
        hits.append({
            "text": doc,
            "source": (meta or {}).get("source", "?"),
            "page": (meta or {}).get("page", None),
            "chunk": (meta or {}).get("chunk", 0),
            "score": round(1.0 - float(dist), 3),
        })
    return hits


def list_sources() -> list[str]:
    coll = _get_collection()
    if coll.count() == 0:
        return []
    data = coll.get(include=["metadatas"])
    metas = data.get("metadatas") or []
    return sorted({m["source"] for m in metas if m and "source" in m})


def clear_index() -> None:
    global _client, _collection
    if _client is None:
        _client = chromadb.PersistentClient(
            path=str(CHROMA_DIR),
            settings=Settings(anonymized_telemetry=False),
        )
    try:
        _client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    _collection = None
