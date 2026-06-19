"""Build a McKinsey-style .pptx deck from a structured spec.

Style:
- White background, navy action title (serif), thin blue accent line under it,
  optional gray lead-in, body content, source footer + page number.
- One idea per slide. Sparse layouts.

The agent calls store_deck(spec) via the `generate_deck` tool; the resulting
bytes are kept in an in-memory dict so the Streamlit UI can offer a download
button without round-tripping the binary through the LLM.
"""
from __future__ import annotations

import io
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Palette — navy / accent blue / black / gray on white
NAVY     = RGBColor(0x0A, 0x25, 0x40)   # action titles
ACCENT   = RGBColor(0x2E, 0x7D, 0xDD)   # thin rule + chart bars + table header
BLACK    = RGBColor(0x1A, 0x1A, 0x1A)   # body text
GRAY     = RGBColor(0x6E, 0x6E, 0x6E)   # lead-in, source, page numbers
LIGHT    = RGBColor(0xE8, 0xEC, 0xF1)   # alt row shading
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Georgia"     # serif for titles
BODY_FONT  = "Calibri"     # sans-serif for body

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_DECKS: dict[str, dict] = {}


# ---------- low-level text helpers ----------

def _style_run(run, *, text: str, font: str, size: int,
               bold: bool = False, italic: bool = False, color=BLACK) -> None:
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _set_text(frame, text: str, *, font: str = BODY_FONT, size: int = 14,
              bold: bool = False, italic: bool = False, color=BLACK,
              align=PP_ALIGN.LEFT) -> None:
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _style_run(run, text=text, font=font, size=size, bold=bold, italic=italic, color=color)


def _add_bullets(frame, bullets: list[str], *, size: int = 16) -> None:
    frame.clear()
    frame.word_wrap = True
    for i, b in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        _style_run(run, text=f"—  {b}", font=BODY_FONT, size=size, color=BLACK)


# ---------- shared slide chrome ----------

def _add_action_title(slide, title: str, lead_in: str | None) -> Emu:
    """Draw the action title + optional lead-in + thin accent rule. Returns the Y where body should start."""
    # Action title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75))
    _set_text(tb.text_frame, title, font=TITLE_FONT, size=22, bold=True, color=NAVY)

    body_top = Inches(1.25)

    # Optional lead-in
    if lead_in:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5))
        _set_text(sb.text_frame, lead_in, font=BODY_FONT, size=13, italic=True, color=GRAY)
        body_top = Inches(1.75)

    # Thin accent rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), body_top - Inches(0.1),
        Inches(1.2), Emu(20000)
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT

    return body_top


def _add_footer(slide, source: str | None, page_no: int) -> None:
    # Source (bottom-left, gray italic, small)
    src_text = f"Source: {source}" if source else ""
    if src_text:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(10.5), Inches(0.3))
        _set_text(sb.text_frame, src_text, font=BODY_FONT, size=9, italic=True, color=GRAY)
    # Page number (bottom-right)
    pb = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.6), Inches(0.3))
    _set_text(pb.text_frame, str(page_no), font=BODY_FONT, size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ---------- slide builders ----------

def _add_cover_slide(prs, title: str, subtitle: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # White background is default; add a thin top navy band for identity
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY

    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(11.9), Inches(1.6))
    _set_text(tb.text_frame, title, font=TITLE_FONT, size=36, bold=True, color=NAVY)

    # Accent rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.0), Inches(1.6), Emu(28000)
    )
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT

    # Subtitle
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.6))
        _set_text(sb.text_frame, subtitle, font=BODY_FONT, size=18, color=GRAY)


def _add_bullets_slide(prs, title: str, lead_in: str | None,
                       bullets: list[str], source: str | None, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body_top = _add_action_title(slide, title, lead_in)
    body = slide.shapes.add_textbox(Inches(0.6), body_top, Inches(12.1), Inches(5.0))
    _add_bullets(body.text_frame, bullets or ["(no content)"])
    _add_footer(slide, source, page_no)


def _add_table_slide(prs, title: str, lead_in: str | None,
                     headers: list[str], rows: list[list[str]],
                     source: str | None, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body_top = _add_action_title(slide, title, lead_in)
    if not headers:
        _add_footer(slide, source, page_no)
        return
    n_rows = max(len(rows) + 1, 2)
    n_cols = len(headers)
    left, top = Inches(0.5), body_top
    width = Inches(12.3)
    row_h = 0.40
    height = Inches(min(row_h * n_rows + 0.3, 5.4))
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_text(cell.text_frame, h, font=BODY_FONT, size=12, bold=True, color=WHITE)
    # Body rows with zebra striping
    for i, row in enumerate(rows, start=1):
        for j in range(n_cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
            val = row[j] if j < len(row) else ""
            _set_text(cell.text_frame, str(val), font=BODY_FONT, size=11, color=BLACK)
    _add_footer(slide, source, page_no)


_CHART_MAP = {
    "bar":    XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "hbar":   XL_CHART_TYPE.BAR_CLUSTERED,
    "line":   XL_CHART_TYPE.LINE,
    "pie":    XL_CHART_TYPE.PIE,
}


def _add_chart_slide(prs, title: str, lead_in: str | None,
                     categories: list[str], series: list[dict],
                     chart_type: str, source: str | None, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    body_top = _add_action_title(slide, title, lead_in)
    if not categories or not series:
        tb = slide.shapes.add_textbox(Inches(0.6), body_top, Inches(12), Inches(4))
        _set_text(tb.text_frame, "(chart data missing)", font=BODY_FONT, size=14, color=GRAY)
        _add_footer(slide, source, page_no)
        return
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        vals = tuple(float(v) for v in s.get("values", []))
        cd.add_series(s.get("name", "Series"), vals)
    xl = _CHART_MAP.get((chart_type or "bar").lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart = slide.shapes.add_chart(
        xl, Inches(0.6), body_top, Inches(12.1), Inches(5.0), cd
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    _add_footer(slide, source, page_no)


# ---------- public API ----------

def build_deck(spec: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_cover_slide(prs, spec.get("title") or "Untitled Deck", spec.get("subtitle"))

    for idx, s in enumerate(spec.get("slides", []), start=2):
        kind = (s.get("type") or "bullets").lower()
        stitle  = s.get("title", "")
        lead_in = s.get("lead_in") or s.get("subtitle")
        source  = s.get("source")
        if kind == "title":
            _add_cover_slide(prs, stitle, s.get("subtitle"))
        elif kind == "table":
            _add_table_slide(prs, stitle, lead_in,
                             s.get("headers", []), s.get("rows", []),
                             source, idx)
        elif kind == "chart":
            _add_chart_slide(prs, stitle, lead_in,
                             s.get("categories", []), s.get("series", []),
                             s.get("chart_type", "bar"), source, idx)
        else:
            _add_bullets_slide(prs, stitle, lead_in, s.get("bullets", []), source, idx)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def store_deck(spec: dict) -> dict:
    """Build a deck and stash it in memory. Returns a small handle for the LLM."""
    data = build_deck(spec)
    deck_id = uuid.uuid4().hex[:12]
    filename = (spec.get("filename") or spec.get("title") or "deck").strip()
    if not filename.lower().endswith(".pptx"):
        filename = f"{filename}.pptx"
    _DECKS[deck_id] = {"bytes": data, "filename": filename}
    return {
        "deck_id": deck_id,
        "filename": filename,
        "size_bytes": len(data),
        "slides_count": len(spec.get("slides", [])) + 1,
        "note": "Deck generated. The UI will render a download button.",
    }


def get_deck(deck_id: str) -> dict | None:
    return _DECKS.get(deck_id)
